"""
Document ingestion pipeline for the Evidence Checker Agent.

Extracts text from uploaded attachments, chunks it, and upserts into a
task-scoped Qdrant collection for semantic search during evidence evaluation.

Supported types:
  - PDF         → PyMuPDF (fitz) page-by-page text + image vision description
  - DOCX        → python-docx heading-hierarchy chunks
  - XLSX / CSV  → openpyxl / pandas per-sheet tabular summary
  - PPTX        → python-pptx per-slide chunks
  - Images      → LiteLLM vision model description
  - Fallback    → raw UTF-8 decode attempt

Design invariants:
  - Idempotent: chunk point IDs are sha256(attachment_id + str(chunk_index))
    so re-running the same attachment is safe (upsert overwrites).
  - Page cap: stop after page_cap pages to prevent runaway costs.
  - All exceptions caught per-attachment; one bad file does not abort the job.
"""

from __future__ import annotations

import hashlib
import io
import json
import os
import time
import uuid
from importlib import import_module
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    import asyncpg

from .models import ChunkResult, IngestionResult

_logging_module = import_module("backend.01_core.logging_utils")
_logger = _logging_module.get_logger("backend.ai.evidence_checker.ingestion")

# ── Configurable knobs ──────────────────────────────────────────────────────
_DEFAULT_PAGE_CAP = int(os.getenv("EVIDENCE_CHECKER_PAGE_CAP", "10000"))
_CHUNK_SIZE = int(os.getenv("EVIDENCE_CHECKER_CHUNK_TOKENS", "512"))
_CHUNK_OVERLAP = int(os.getenv("EVIDENCE_CHECKER_CHUNK_OVERLAP", "64"))
_QDRANT_COLLECTION_PREFIX = "evcheck_task_"

# ── MIME type routing ───────────────────────────────────────────────────────
_PDF_MIMES = {"application/pdf"}
_DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_XLSX_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/csv",
}
_PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
_IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/tiff", "image/bmp", "image/webp"}


# ── Chunk helpers ────────────────────────────────────────────────────────────

def _point_id(attachment_id: str, chunk_index: int) -> str:
    """Stable deterministic UUID from attachment + chunk index."""
    raw = f"{attachment_id}:{chunk_index}".encode()
    digest = hashlib.sha256(raw).hexdigest()
    return str(uuid.UUID(digest[:32]))


def _simple_chunk(text: str, chunk_size: int = _CHUNK_SIZE, overlap: int = _CHUNK_OVERLAP) -> list[str]:
    """
    Word-boundary chunker.  Splits on whitespace, groups words into chunks of
    ~chunk_size words, with overlap.  No external tokenizer dependency.
    """
    words = text.split()
    if not words:
        return []
    step = max(1, chunk_size - overlap)
    chunks = []
    i = 0
    while i < len(words):
        chunk_words = words[i: i + chunk_size]
        chunks.append(" ".join(chunk_words))
        i += step
    return chunks


# ── Extractors ───────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes, page_cap: int) -> list[dict]:
    """Returns list of {page_number, text, section_or_sheet}."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        _logger.warning("PyMuPDF not installed; PDF extraction skipped")
        return []

    pages = []
    doc = fitz.open(stream=data, filetype="pdf")
    for i, page in enumerate(doc):
        if i >= page_cap:
            break
        text = page.get_text("text") or ""
        pages.append({"page_number": i + 1, "text": text.strip(), "section_or_sheet": None})
    doc.close()
    return pages


def _extract_docx(data: bytes) -> list[dict]:
    try:
        from docx import Document
    except ImportError:
        _logger.warning("python-docx not installed; DOCX extraction skipped")
        return []

    doc = Document(io.BytesIO(data))
    sections: list[dict] = []
    current_heading = "Document"
    current_paras: list[str] = []

    def flush():
        if current_paras:
            text = " ".join(current_paras)
            sections.append({
                "page_number": None,
                "text": text,
                "section_or_sheet": current_heading,
            })

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            flush()
            current_heading = para.text.strip() or current_heading
            current_paras = []
        else:
            stripped = para.text.strip()
            if stripped:
                current_paras.append(stripped)

    flush()
    return sections


def _extract_xlsx(data: bytes) -> list[dict]:
    try:
        import openpyxl
    except ImportError:
        _logger.warning("openpyxl not installed; XLSX extraction skipped")
        return []

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    pages = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):
                rows_text.append(" | ".join(cells))
            if len(rows_text) > 2000:  # cap per sheet
                rows_text.append("... (sheet truncated)")
                break
        if rows_text:
            pages.append({
                "page_number": None,
                "text": "\n".join(rows_text),
                "section_or_sheet": sheet_name,
            })
    wb.close()
    return pages


def _extract_csv(data: bytes) -> list[dict]:
    text = data.decode("utf-8", errors="replace")
    return [{"page_number": None, "text": text[:50000], "section_or_sheet": "Sheet1"}]


def _extract_pptx(data: bytes) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        _logger.warning("python-pptx not installed; PPTX extraction skipped")
        return []

    prs = Presentation(io.BytesIO(data))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({
                "page_number": i,
                "text": " ".join(texts),
                "section_or_sheet": None,
            })
    return pages


def _extract_fallback(data: bytes) -> list[dict]:
    try:
        text = data.decode("utf-8", errors="replace").strip()
        if text:
            return [{"page_number": None, "text": text[:100000], "section_or_sheet": None}]
    except Exception:
        pass
    return []


async def _describe_image_via_llm(
    data: bytes,
    filename: str,
    provider,  # LLMProvider with vision support
) -> str:
    """Ask the LLM to describe the image content for indexing."""
    import base64
    b64 = base64.b64encode(data).decode()
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "png"
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "tiff": "image/tiff", "bmp": "image/bmp", "webp": "image/webp"}
    mime = mime_map.get(ext, "image/png")

    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime};base64,{b64}"},
                },
                {
                    "type": "text",
                    "text": (
                        "Describe all text, data, and compliance-relevant content visible in this image. "
                        "Include any signatures, dates, policy references, control names, or attestation statements. "
                        "Be thorough — this description will be used to verify compliance evidence."
                    ),
                },
            ],
        }
    ]
    try:
        response = await provider.chat_completion(
            messages=messages, tools=None, temperature=1.0, max_tokens=1024
        )
        return response.content or "[image: no description produced]"
    except Exception as exc:
        _logger.warning("Image description failed for %s: %s", filename, exc)
        return f"[image: description failed — {exc}]"


# ── Qdrant indexer ───────────────────────────────────────────────────────────

def _collection_name(task_id: str) -> str:
    # Use a single global collection for all evidence checker documents
    return "evidence_checker_docs"


async def _ensure_collection(qdrant_client, name: str, vector_size: int) -> None:
    from qdrant_client.models import Distance, VectorParams
    try:
        qdrant_client.get_collection(name)
    except Exception:
        qdrant_client.create_collection(
            collection_name=name,
            vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
        )


def _get_embed_provider(chat_provider):
    """
    Return the provider to use for embeddings.

    Priority:
    1. EVIDENCE_CHECKER_EMBED_URL + EVIDENCE_CHECKER_EMBED_KEY (explicit override)
    2. AI_EMBEDDING_URL + AI_EMBEDDING_API_KEY (shared embedding endpoint)
    3. The chat provider itself if it has an `embed` method (OpenAI / Azure OpenAI)
    4. Raises RuntimeError if no embedding capability is available
    """
    _llm_factory = import_module("backend.20_ai.14_llm_providers.factory")

    # Option 1: explicit evidence checker embedding endpoint
    embed_url = os.getenv("EVIDENCE_CHECKER_EMBED_URL", "").strip()
    embed_key = os.getenv("EVIDENCE_CHECKER_EMBED_KEY", "").strip()
    embed_model = os.getenv("EVIDENCE_CHECKER_EMBED_MODEL", "").strip()

    if embed_url and embed_key:
        return _llm_factory.get_provider(
            provider_type="openai_compatible",
            provider_base_url=embed_url,
            api_key=embed_key,
            model_id=embed_model or "text-embedding-3-small",
            temperature=1.0,
        )

    # Option 2: shared AI embedding endpoint (AI_EMBEDDING_URL / AI_EMBEDDING_API_KEY)
    ai_embed_url = os.getenv("AI_EMBEDDING_URL", "").strip()
    ai_embed_key = os.getenv("AI_EMBEDDING_API_KEY", "").strip()
    ai_embed_model = os.getenv("AI_EMBEDDING_MODEL", "text-embedding-3-small").strip()

    if ai_embed_url and ai_embed_key:
        return _llm_factory.get_provider(
            provider_type="openai_compatible",
            provider_base_url=ai_embed_url,
            api_key=ai_embed_key,
            model_id=ai_embed_model or "text-embedding-3-small",
            temperature=1.0,
        )

    # Option 3: use chat provider's own embed() method
    if hasattr(chat_provider, "embed") and callable(chat_provider.embed):
        return chat_provider

    raise RuntimeError(
        "No embedding provider available. Set AI_EMBEDDING_URL + AI_EMBEDDING_API_KEY "
        "env vars, or use an OpenAI/Azure OpenAI chat provider."
    )


async def _embed_texts(texts: list[str], provider) -> list[list[float]]:
    """Use the resolved embedding provider. Batched in groups of 64."""
    embed_provider = _get_embed_provider(provider)
    embeddings: list[list[float]] = []
    batch_size = 64
    for i in range(0, len(texts), batch_size):
        batch = texts[i: i + batch_size]
        try:
            result = await embed_provider.embed(batch)
            embeddings.extend(result)
        except Exception as exc:
            _logger.error("Embedding batch %d failed: %s", i // batch_size, exc)
            # Zero vectors are stored but won't surface in semantic search —
            # the full corpus map phase (get_all_chunks) will still scan them.
            for _ in batch:
                embeddings.append([0.0] * 1536)
    return embeddings


async def index_attachment(
    *,
    task_id: str,
    org_id: str,
    attachment_id: str,
    document_filename: str,
    file_data: bytes,
    mime_type: str,
    page_cap: int,
    llm_provider,  # LLMProvider — used for vision + embedding
    qdrant_client,
) -> IngestionResult:
    """
    Extract → chunk → embed → upsert one attachment into the task's Qdrant collection.
    Idempotent: same chunk_index produces same point ID (upsert).
    """
    t0 = time.monotonic()
    mime = (mime_type or "").lower().split(";")[0].strip()

    # ── Extract raw pages/sections ──────────────────────────────────────────
    if mime in _PDF_MIMES:
        raw_pages = _extract_pdf(file_data, page_cap)
    elif mime in _DOCX_MIMES:
        raw_pages = _extract_docx(file_data)
    elif mime in _XLSX_MIMES:
        raw_pages = _extract_xlsx(file_data) if "csv" not in mime else _extract_csv(file_data)
    elif mime in _PPTX_MIMES:
        raw_pages = _extract_pptx(file_data)
    elif mime in _IMAGE_MIMES:
        description = await _describe_image_via_llm(file_data, document_filename, llm_provider)
        raw_pages = [{"page_number": None, "text": description, "section_or_sheet": None}]
    else:
        raw_pages = _extract_fallback(file_data)

    if not raw_pages:
        return IngestionResult(
            attachment_id=attachment_id,
            pages_processed=0,
            chunks_indexed=0,
            error="no extractable content",
        )

    pages_processed = len(raw_pages)

    # ── Chunk ────────────────────────────────────────────────────────────────
    all_chunks: list[dict] = []
    chunk_index = 0
    for page in raw_pages:
        text = page["text"]
        if not text:
            continue
        for chunk_text in _simple_chunk(text):
            all_chunks.append({
                "chunk_index": chunk_index,
                "text": chunk_text,
                "page_number": page["page_number"],
                "section_or_sheet": page["section_or_sheet"],
            })
            chunk_index += 1

    if not all_chunks:
        return IngestionResult(
            attachment_id=attachment_id,
            pages_processed=pages_processed,
            chunks_indexed=0,
            error="content extracted but no chunks produced",
        )

    # ── Embed ────────────────────────────────────────────────────────────────
    texts = [c["text"] for c in all_chunks]
    embeddings = await _embed_texts(texts, llm_provider)
    vector_size = len(embeddings[0]) if embeddings and embeddings[0] else 1536

    # ── Upsert to Qdrant ─────────────────────────────────────────────────────
    collection = _collection_name(task_id)
    await _ensure_collection(qdrant_client, collection, vector_size)

    from qdrant_client.models import PointStruct
    points = []
    for chunk, vector in zip(all_chunks, embeddings):
        points.append(PointStruct(
            id=_point_id(attachment_id, chunk["chunk_index"]),
            vector=vector,
            payload={
                "org_id": org_id,
                "task_id": task_id,
                "attachment_id": attachment_id,
                "document_filename": document_filename,
                "page_number": chunk["page_number"],
                "section_or_sheet": chunk["section_or_sheet"],
                "chunk_index": chunk["chunk_index"],
                "text": chunk["text"],
            },
        ))

    # Upsert in batches of 100
    for i in range(0, len(points), 100):
        qdrant_client.upsert(collection_name=collection, points=points[i: i + 100])

    elapsed = round(time.monotonic() - t0, 2)
    _logger.info(
        "evidence_check.ingestion_done",
        extra={
            "task_id": task_id,
            "attachment_id": attachment_id,
            "pages": pages_processed,
            "chunks": len(all_chunks),
            "elapsed_s": elapsed,
        },
    )
    return IngestionResult(
        attachment_id=attachment_id,
        pages_processed=pages_processed,
        chunks_indexed=len(all_chunks),
    )


async def search_collection(
    *,
    task_id: str,
    org_id: str,
    query: str,
    top_k: int,
    llm_provider,
    qdrant_client,
) -> list[ChunkResult]:
    """Semantic search in the task's Qdrant collection, filtered to this org."""
    collection = _collection_name(task_id)

    # Embed the query
    try:
        vectors = await _embed_texts([query], llm_provider)
        query_vector = vectors[0]
    except Exception as exc:
        _logger.error("Query embedding failed: %s", exc)
        return []

    from qdrant_client.models import Filter, FieldCondition, MatchValue
    try:
        results = qdrant_client.search(
            collection_name=collection,
            query_vector=query_vector,
            limit=top_k,
            query_filter=Filter(
                must=[
                    FieldCondition(key="org_id", match=MatchValue(value=org_id)),
                    FieldCondition(key="task_id", match=MatchValue(value=task_id))
                ]
            ),
            with_payload=True,
        )
    except Exception as exc:
        _logger.warning("Qdrant search failed for task %s: %s", task_id, exc)
        return []

    chunks = []
    for hit in results:
        p = hit.payload or {}
        chunks.append(ChunkResult(
            text=p.get("text", ""),
            document_filename=p.get("document_filename", "unknown"),
            page_number=p.get("page_number"),
            section_or_sheet=p.get("section_or_sheet"),
            attachment_id=p.get("attachment_id", ""),
            chunk_index=p.get("chunk_index", 0),
            score=hit.score,
        ))
    return chunks


def delete_collection(*, task_id: str, qdrant_client) -> None:
    """Delete all chunks for a specific task_id from the global collection."""
    collection = _collection_name(task_id)
    from qdrant_client.models import Filter, FieldCondition, MatchValue
    try:
        qdrant_client.delete(
            collection_name=collection,
            points_selector=Filter(
                must=[FieldCondition(key="task_id", match=MatchValue(value=task_id))]
            ),
        )
        _logger.info("evidence_check.task_chunks_deleted", extra={"task_id": task_id})
    except Exception as exc:
        _logger.warning("Could not delete chunks for task %s: %s", task_id, exc)


async def get_all_chunks(
    *,
    task_id: str,
    org_id: str,
    qdrant_client,
    batch_size: int = 500,
) -> list[ChunkResult]:
    """
    Retrieve ALL indexed chunks for a task (not just top-K).
    Used by the map-reduce evaluator to guarantee full corpus coverage.
    Scrolls through the Qdrant collection in batches.
    """
    collection = _collection_name(task_id)
    from qdrant_client.models import Filter, FieldCondition, MatchValue

    task_filter = Filter(
        must=[
            FieldCondition(key="org_id", match=MatchValue(value=org_id)),
            FieldCondition(key="task_id", match=MatchValue(value=task_id)),
        ]
    )

    all_chunks: list[ChunkResult] = []
    offset = None

    try:
        while True:
            result = qdrant_client.scroll(
                collection_name=collection,
                scroll_filter=task_filter,
                limit=batch_size,
                offset=offset,
                with_payload=True,
                with_vectors=False,
            )
            points, next_offset = result
            for point in points:
                p = point.payload or {}
                all_chunks.append(ChunkResult(
                    text=p.get("text", ""),
                    document_filename=p.get("document_filename", "unknown"),
                    page_number=p.get("page_number"),
                    section_or_sheet=p.get("section_or_sheet"),
                    attachment_id=p.get("attachment_id", ""),
                    chunk_index=p.get("chunk_index", 0),
                    score=1.0,  # not scored — full scan
                ))
            if next_offset is None:
                break
            offset = next_offset
    except Exception as exc:
        _logger.warning("get_all_chunks failed for task %s: %s", task_id, exc)
        return []

    return all_chunks


def delete_attachment_chunks(*, task_id: str, attachment_id: str, qdrant_client) -> None:
    """Delete all chunks for a specific attachment from the global collection."""
    collection = _collection_name(task_id)
    from qdrant_client.models import Filter, FieldCondition, MatchValue
    try:
        qdrant_client.delete(
            collection_name=collection,
            points_selector=Filter(
                must=[
                    FieldCondition(key="task_id", match=MatchValue(value=task_id)),
                    FieldCondition(key="attachment_id", match=MatchValue(value=attachment_id)),
                ]
            ),
        )
        _logger.info("evidence_check.attachment_chunks_deleted", extra={"task_id": task_id, "attachment_id": attachment_id})
    except Exception as exc:
        _logger.warning("Could not delete chunks for task %s, attachment %s: %s", task_id, attachment_id, exc)
